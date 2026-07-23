from pptx import Presentation


def parse(file_path):

    prs = Presentation(file_path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    text = text.strip()

    if not text:
        raise Exception("Empty presentation")

    return text